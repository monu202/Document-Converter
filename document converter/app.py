from flask import Flask, request, send_file, render_template
from werkzeug.utils import secure_filename
import os
import fitz  # PyMuPDF
from docx import Document
from docx.shared import Pt
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx2pdf import convert
from pptx import Presentation
from reportlab.pdfgen import canvas
from pdf2docx import Converter

app = Flask(__name__, template_folder='templates')

# Configuration for upload and converted folders
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['CONVERTED_FOLDER'] = 'converted'
app.config['ALLOWED_EXTENSIONS'] = {'pdf', 'docx', 'pptx'}

# Ensure folders exist
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['CONVERTED_FOLDER'], exist_ok=True)

# Function to check allowed file extensions
def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']

@app.route('/', methods=['GET'])
def index():
    return render_template('base.html')

@app.route('/convert', methods=['POST'])
def convert_file():
    if 'file' not in request.files:
        return "No file part in the request.", 400

    file = request.files['file']
    if file.filename == '':
        return "No file selected.", 400

    if not allowed_file(file.filename):
        return "Invalid file type selected. Please upload a PDF, Word, or PowerPoint file.", 400

    conversion_type = request.form['conversionType']
    filename = secure_filename(file.filename)
    filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(filepath)

    if conversion_type == 'pdf-to-word':
        output_path = convert_pdf_to_word(filepath, filename)
    elif conversion_type == 'word-to-pdf':
        output_path = convert_word_to_pdf(filepath, filename)
    elif conversion_type == 'ppt-to-pdf':
        output_path = convert_ppt_to_pdf(filepath, filename)
    elif conversion_type == 'pdf-to-ppt':
        output_path = convert_pdf_to_ppt(filepath, filename)
    else:
        return "Invalid conversion type selected.", 400

    return send_file(output_path, as_attachment=True)

def convert_pdf_to_word(filepath, filename):
    output_path = os.path.join(app.config['CONVERTED_FOLDER'], filename.replace('.pdf', '.docx'))
    converter = Converter(filepath)  # Initialize the converter
    try:
        converter.convert(output_path, start=0, end=None)  # Perform the conversion
    finally:
        converter.close()  # Ensure the converter is closed after use
    return output_path

def convert_word_to_pdf(filepath, filename):
    # Convert Word to PDF preserving formatting using docx2pdf
    output_path = os.path.join(app.config['CONVERTED_FOLDER'], filename.replace('.docx', '.pdf'))
    try:
        convert(filepath, output_path)
    except Exception as e:
        return f"Error converting Word to PDF: {e}", 500
    return output_path

def convert_ppt_to_pdf(filepath, filename):
    # Convert PowerPoint to PDF using python-pptx and reportlab
    output_path = os.path.join(app.config['CONVERTED_FOLDER'], filename.replace('.pptx', '.pdf'))
    ppt = Presentation(filepath)
    c = canvas.Canvas(output_path)

    for slide in ppt.slides:
        y_position = 750
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                c.drawString(50, y_position, text)
                y_position -= 20  # Move down for the next line of text

    c.save()
    return output_path

def convert_pdf_to_ppt(filepath, filename):
    output_path = os.path.join(app.config['CONVERTED_FOLDER'], filename.replace('.pdf', '.pptx'))
    ppt = Presentation()
    pdf_reader = fitz.open(filepath)

    for page_num, page in enumerate(pdf_reader):
        slide = ppt.slides.add_slide(ppt.slide_layouts[5])  # Blank slide layout

        # Add text
        text_box = slide.shapes.add_textbox(left=Pt(50), top=Pt(50), width=Pt(500), height=Pt(400))
        text_frame = text_box.text_frame

        blocks = page.get_text("dict")["blocks"]
        for block in blocks:
            if "lines" in block:
                for line in block["lines"]:
                    paragraph = text_frame.add_paragraph()
                    paragraph.text = " ".join(span["text"] for span in line["spans"])

        # Add images
        images = page.get_images(full=True)
        for img_index, img in enumerate(images):
            xref = img[0]
            base_image = pdf_reader.extract_image(xref)
            image_bytes = base_image["image"]
            image_filename = f"temp_image_{page_num}_{img_index}.png"
            image_path = os.path.join(app.config['UPLOAD_FOLDER'], image_filename)

            with open(image_path, "wb") as img_file:
                img_file.write(image_bytes)

            slide.shapes.add_picture(image_path, left=Pt(50), top=Pt(450))

    ppt.save(output_path)
    return output_path

if __name__ == '__main__':
    app.run(debug=True)
